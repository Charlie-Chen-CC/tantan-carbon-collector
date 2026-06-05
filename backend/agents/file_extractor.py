"""
文件提取Agent - 碳管师收资系统
负责从用户上传的各类文件中提取对应部分的数据
支持格式: xlsx, xls, pdf, docx, doc, pptx, md
每个Section使用专精的LLM提示词进行数据提取
"""

import uuid
import os
import base64
from datetime import datetime
from typing import Dict, Any, Optional, List
from io import BytesIO

import openpyxl
import docx
import pptx
import pdfplumber
from PIL import Image

from tantan.backend.utils import get_logger
from tantan.backend.rag import get_llm_client
from tantan.backend.agents.pdf_splitter import PDFSplitter

logger = get_logger(__name__)


# ============== 通用文本提取器 ==============

class TextExtractor:
    """从各类文件提取原始文本"""
    _pdf_splitter = PDFSplitter(page_split_threshold=5)

    @staticmethod
    def extract_from_bytes(file_content: bytes, filename: str) -> str:
        """根据文件扩展名调用对应提取器"""
        _, ext = os.path.splitext(filename or '')
        ext = ext.lower()

        extractors = {
            '.xlsx': TextExtractor._extract_xlsx,
            '.xls': TextExtractor._extract_xls,
            '.pdf': TextExtractor._extract_pdf,
            '.docx': TextExtractor._extract_docx,
            '.doc': TextExtractor._extract_doc,
            '.pptx': TextExtractor._extract_pptx,
            '.md': TextExtractor._extract_md,
            '.png': TextExtractor._extract_image,
            '.jpg': TextExtractor._extract_image,
            '.jpeg': TextExtractor._extract_image,
        }

        extractor = extractors.get(ext)
        if not extractor:
            raise ValueError(f"不支持的文件类型: {ext}")

        return extractor(file_content)

    @staticmethod
    def _extract_xlsx(file_content: bytes) -> str:
        """从 xlsx 提取文本"""
        try:
            wb = openpyxl.load_workbook(BytesIO(file_content), data_only=True)
            texts = []
            for sheet in wb.worksheets:
                texts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            texts.append(str(cell.value))
            return "\n".join(texts)
        except Exception as e:
            logger.error(f"xlsx 提取失败: {e}")
            return ""

    @staticmethod
    def _extract_xls(file_content: bytes) -> str:
        """从 xls 提取文本（复用 xlsx 逻辑）"""
        return TextExtractor._extract_xlsx(file_content)

    @staticmethod
    def _extract_pdf(file_content: bytes) -> str:
        """从 pdf 提取文本（支持分页分割）"""
        try:
            segments = TextExtractor._pdf_splitter.split(file_content)
            return "\n\n".join(segments)
        except Exception as e:
            logger.error(f"pdf 提取失败: {e}")
            return ""

    @staticmethod
    def _extract_docx(file_content: bytes) -> str:
        """从 docx 提取文本"""
        try:
            doc = docx.Document(BytesIO(file_content))
            texts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    texts.append(para.text)
            # 提取表格
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        texts.append(" | ".join(row_text))
            return "\n".join(texts)
        except Exception as e:
            logger.error(f"docx 提取失败: {e}")
            return ""

    @staticmethod
    def _extract_doc(file_content: bytes) -> str:
        """从 doc 提取文本（尝试直接读取二进制）"""
        # python-docx 不支持 .doc，但可以直接读二进制文本
        try:
            text = file_content.decode("utf-8", errors="ignore")
            # 简单过滤不可见字符
            import re
            text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f]', '', text)
            return text
        except Exception as e:
            logger.error(f"doc 提取失败: {e}")
            return ""

    @staticmethod
    def _extract_pptx(file_content: bytes) -> str:
        """从 pptx 提取文本"""
        try:
            prs = pptx.Presentation(BytesIO(file_content))
            texts = []
            for i, slide in enumerate(prs.slides):
                texts.append(f"[Slide {i+1}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception as e:
            logger.error(f"pptx 提取失败: {e}")
            return ""

    @staticmethod
    def _extract_md(file_content: bytes) -> str:
        """从 markdown 提取文本

        自动探测编码：优先尝试 UTF-8（带/不带 BOM），失败则尝试 GBK（Windows 常见），
        再不行降级到 latin-1（永不抛异常）。避免硬编码 UTF-8 把中文 MD 解析成乱码。
        """
        for encoding in ("utf-8-sig", "utf-8", "gbk", "gb18030", "latin-1"):
            try:
                return file_content.decode(encoding)
            except (UnicodeDecodeError, LookupError):
                continue
        # 兜底：用 errors="ignore" 不抛异常（跟原行为一致）
        return file_content.decode("utf-8", errors="ignore")

    @staticmethod
    def _extract_image(file_content: bytes) -> str:
        """从图片提取文本：返回base64编码供LLM处理"""
        try:
            image = Image.open(BytesIO(file_content))
            width, height = image.size

            # 检测图片尺寸，过小的图片可能是截图或图标
            if width < 300 or height < 300:
                return "[图片尺寸过小，跳过]"

            # 将图片内容转为base64供多模态LLM使用
            b64_content = base64.b64encode(file_content).decode('utf-8')

            # 判断图片格式
            img_format = image.format.lower() if hasattr(image, 'format') else 'jpeg'
            mime_type = f"image/{img_format}" if img_format != 'jpeg' else "image/jpeg"

            return f"[IMAGE_DATA]{b64_content}[/IMAGE_DATA] mime:{mime_type} size:{width}x{height}"
        except Exception as e:
            logger.error(f"图片 提取失败: {e}")
            return ""


# ============== Section 专精提示词 ==============

from tantan.backend.agents.section_options import (
    INDUSTRY_STANDARD_REF,
    INDUSTRY_EXAMPLES,
    PRODUCT_UNITS,
    PRODUCT_OTHER_MAX,
    COMMON_FUEL_TYPES,
    MEASURABLE_OPTIONS,
    PHOTOVOLTAIC_CONFIG_OPTIONS,
    GREEN_CERTIFICATE_OPTIONS,
    EMISSION_RIGHTS_OPTIONS,
    COMMON_REFRIGERANTS,
    FIRE_EXTINGUISHER_TYPES,
    WASTE_DISPOSAL_METHODS,
    TRANSPORT_MODES,
    STATISTICAL_CALIBER_OPTIONS,
)


def _join_or(options: list[str]) -> str:
    """把选项列表拼成"a/b/c"格式"""
    return "、".join(options)


def _build_section_prompts() -> dict[int, str]:
    """构建 9 个 Section 的 LLM 提示词（Phase 6.1 反硬编码）

    所有选项枚举引用 `section_options.py`（后续 Phase 6.2 会迁到 `shared/field_schema.json` + codegen 注入）。
    提示词字符串里不再出现"燃料类型包括：天然气、液化石油气、煤、柴油、汽油"
    这类硬编码列表，统一改为"请参照 {OPTIONS} 理解业务背景，但不应限制 LLM 自行识别其他合理取值"。
    """
    return {
        1: f"""你是一个碳排放数据填写助手，专门从【碳盘查企业碳排放核算报告】或【数据收集表】中提取【基本信息】部分的数据。

【业务背景说明】
- 基本信息是企业进行碳排放核算的基础
- 企业名称通常是公司全称，如"浙江光华科技股份有限公司"
- 所属行业按{INDUSTRY_STANDARD_REF}填写
- 核算年份指碳排放核算的时间范围，一般为一年（1月1日至12月31日）

【字段示例（仅供理解，**不**是穷举）】
- 所属行业示例：{_join_or(INDUSTRY_EXAMPLES)}（不同企业可能涉及其他行业分类）

需要提取的字段：
- 企业名称（公司全称）
- 所属行业（按{INDUSTRY_STANDARD_REF}填写）
- 联系人（负责碳排放核算的人员姓名）
- 联系方式（电话或邮箱）
- 生产地址（实际生产场所地址）
- 核算年份（一般为4位数字年份，如2023、2024）
- 核算周期说明（起止日期，如2023.1.1~2023.12.31 或 2024年1月1日至2024年12月31日）

请从提供的原始文本中提取上述字段，以JSON格式返回。字段名为key，值为字符串，如果某字段不存在则返回null。
只返回JSON，不要有其他解释。""",

        2: f"""你是一个碳排放数据填写助手，专门从【碳盘查企业碳排放核算报告】或【数据收集表】中提取【产品】部分的数据。

【业务背景说明】
- PCF (Product Carbon Footprint) 指产品碳足迹
- 核算目标产品是企业进行碳排放核算的主要对象
- 计量单位决定如何量化产品，常见计量单位有：{_join_or(PRODUCT_UNITS)}
- 如果工厂生产多种产品，需要说明是否有副产品产出

需要提取的字段：
- PCF核算目标产品名称（企业主要产品的名称）
- 是否为生产工厂唯一产品（值：是 / 否）
- 其他产品1名称、其他产品2名称...其他产品{PRODUCT_OTHER_MAX}名称（如果不止一种产品）
- 其他产品超过{PRODUCT_OTHER_MAX}种的说明（如超过请说明数量）
- 计量单位
- 目标产品产线内是否有副产品（值：是 / 否）
- 副产品1名称、副产品2名称...副产品{PRODUCT_OTHER_MAX}名称
- 副产品超过{PRODUCT_OTHER_MAX}种的说明

请从提供的原始文本中提取上述字段，以JSON格式返回。只返回JSON。""",

        3: f"""你是一个碳排放数据填写助手，专门从【能源消耗台账】或【碳盘查报告】中提取【燃料使用】部分的数据。

【业务背景说明】
- 燃料使用碳排放是企业碳排放的主要来源之一
- 常见燃料类型有：{_join_or(COMMON_FUEL_TYPES)} 等（**不**是穷举；不同企业可能用生物质成型燃料、煤层气等）
- 生产用锅炉是企业最常见的耗能设备
- 厂内转运车辆（叉车）通常使用柴油或电力
- 商务车和道路车辆燃料指企业自有车辆的用油量

需要提取的字段（每个燃料字段包含：燃料类型、核算周期内使用量、单位、热值、热值单位）：
- 生产用锅炉燃料：{{燃料类型, 使用量, 单位, 实测热值, 热值单位}}
- 专用废气焚烧炉燃料：{{燃料类型, 使用量, 单位}}
- 危废焚烧炉燃料：{{燃料类型, 使用量, 单位}}
- 发电机燃料：{{燃料类型, 使用量, 单位}}
- 食堂炉灶燃料：{{燃料类型, 使用量, 单位}}
- 厂内转运叉车燃料：{{燃料类型, 使用量, 单位}}
- 自有商务车92#：{{燃料类型}}
- 自有商务车95#：{{燃料类型}}
- 自有商务车98#：{{燃料类型}}
- 自有道路车辆燃料-柴油：{{燃料类型}}
- 切割、焊接燃料：{{燃料类型}}

请从提供的原始文本中提取上述字段，转换为JSON格式返回。每个燃料字段返回一个对象，包含上述属性。
只返回JSON，不要有其他解释。""",

        4: f"""你是一个碳排放数据填写助手，专门从【电力消耗台账】或【能源统计表】中提取【电力、热力使用】部分的数据。

【业务背景说明】
- 外购电力是企业碳排放的重要来源（Scope 2）
- 如果企业有光伏发电设施，可以抵扣部分碳排放
- 绿证（绿色电力证书）购买后可抵扣碳排放
- 蒸汽可能为外购或自产，用于生产工艺

【重要】每个字段必须单独返回，不要合并！

需要提取的字段：
- 全厂用电（值：{_join_or(MEASURABLE_OPTIONS)}）
- 生产用电（值：{_join_or(MEASURABLE_OPTIONS)}）
- 行政办公用电（值：{_join_or(MEASURABLE_OPTIONS)}）
- 目标产品产线用电（值：{_join_or(MEASURABLE_OPTIONS)}）
- 单耗用电（值：{_join_or(MEASURABLE_OPTIONS)}）
- 光伏发电量（数字，单位kWh或MWh）
- 光伏发电配置（参考值：{_join_or(PHOTOVOLTAIC_CONFIG_OPTIONS)}）
- 是否购买绿证（参考值：{_join_or(GREEN_CERTIFICATE_OPTIONS)}）
- 是否购买排放权益（参考值：{_join_or(EMISSION_RIGHTS_OPTIONS)}）
- 蒸汽温度（数字，单位℃）
- 蒸汽压力（数字，单位MPa）
- 全厂用蒸汽（值：{_join_or(MEASURABLE_OPTIONS)}）
- 生产用蒸汽（值：{_join_or(MEASURABLE_OPTIONS)}）
- 行政类用蒸汽（值：{_join_or(MEASURABLE_OPTIONS)}）
- 目标产品产线用蒸汽（值：{_join_or(MEASURABLE_OPTIONS)}）
- 单耗用蒸汽（值：{_join_or(MEASURABLE_OPTIONS)}）

请从提供的原始文本中提取上述字段，以JSON格式返回。只返回JSON。""",

        5: f"""你是一个碳排放数据填写助手，专门从【设备清单】或【制冷剂采购发票】中提取【制冷剂使用】部分的数据。

【业务背景说明】
- 制冷剂泄漏是HFC、PFC等强效温室气体的排放源
- 空调和冷冻机是常见的制冷剂使用设备
- 常见制冷剂标号：{_join_or(COMMON_REFRIGERANTS)} 等（**不**是穷举）
- 填充量通常以kg为单位

需要提取的字段：
- 空调制冷剂：数组，包含设备名称、标号、填充量
  例如：[(设备名称: "办公楼1号空调", 标号: "R410A", 填充量: "15.5"), ...]
- 冷冻机制冷剂：数组，包含设备名称、标号、填充量
  例如：[(设备名称: "冷库冷冻机1", 标号: "R22", 填充量: "50"), ...]

请从提供的原始文本中提取上述字段，以JSON格式返回。
只返回JSON，不要有其他解释。""",

        6: f"""你是一个碳排放数据填写助手，专门从【HR人事系统】或【生产报表】中提取【其他散逸类排放】部分的数据。

【业务背景说明】
- CO2灭火器填充属于逸散性排放，填充量即排放量
- 员工工时用于计算焊接、切割等作业的碳排放
- 常见需统计的灭火器类型：{_join_or(FIRE_EXTINGUISHER_TYPES)}（**不**是穷举）

需要提取的字段：
- CO2灭火器填充总量(kg)（灭火器维修填充的CO2重量）
- 核算期内员工总工时(h)（所有员工在核算期内的总工作时间）

请从提供的原始文本中提取上述字段，以JSON格式返回。只返回JSON。""",

        7: f"""你是一个碳排放数据填写助手，专门从【环保设施运行记录】或【废物处理合同】中提取【三废处理】部分的数据。

【业务背景说明】
- 废水处理方式影响废水排放的碳排放因子
- 危废处理常见方式：{_join_or(WASTE_DISPOSAL_METHODS)}（**不**是穷举）
- 污水处理会产生甲烷等温室气体

需要提取的字段：
- 废水处理方式
- 废水处理量（数字，单位t或m3）
- 目标产品产线废水（数字）
- COD浓度（数字，单位mg/L）
- 污水处理药剂1（如果使用）
- 污水处理药剂2（如果使用）
- 污水处理药剂3（如果使用）
- 废气处理方式
- 危废委外焚烧总量（数字）
- 危废委外焚烧目标产品产线分解（数字）
- 危废自行焚烧总量（数字）
- 危废自行焚烧目标产品产线分解（数字）
- 危废委外资源化总量（数字）
- 危废委外资源化目标产品产线分解（数字）
- 危废自行资源化总量（数字）
- 危废自行资源化目标产品产线分解（数字）
- 烟气处理药剂1（如果使用）
- 烟气处理药剂2（如果使用）
- 烟气处理药剂3（如果使用）
- 烟气处理药剂4（如果使用）

请从提供的原始文本中提取上述字段，以JSON格式返回。只返回JSON。""",

        8: f"""你是一个碳排放数据填写助手，专门从【原材料采购清单】或【工艺设计文件】中提取【原材料使用】部分的数据。

【业务背景说明】
- 原材料使用是产品碳足迹的主要贡献者（Scope 3）
- 需要记录所有主要原材料的名称和用量
- 供应商信息用于供应链碳排放核算
- 工艺流程图描述产品生产的完整过程
- 常见运输方式参考：{_join_or(TRANSPORT_MODES)}（**不**是穷举）

需要提取的字段：
- PCF核算目标产品生产工艺流程文字描述
- 原材料：数组，包含名称、规格、使用量、单位
  例如：[(名称: "硫酸", 规格: "98%", 使用量: "100", 单位: "t"), ...]
- 供应商：数组，包含名称、采购的原材料品类、运输方式、运距
  例如：[(名称: "供应商A", 品类: "硫酸", 运输方式: "公路运输", 运距: "200"), ...]

请从提供的原始文本和图片中提取上述字段，以JSON格式返回。
只返回JSON，不要有其他解释。""",

        9: f"""你是一个碳排放数据填写助手，专门从【能源统计表】或【水务账单】中提取【生产耗材】部分的数据。

【业务背景说明】
- 新鲜水消耗是产品碳足迹的一部分
- 氮气常用于化工生产的吹扫和保护
- 统计口径参考：{_join_or(STATISTICAL_CALIBER_OPTIONS)}（**不**是穷举）

需要提取的字段：
- 新鲜水：统计口径（参考值：{_join_or(STATISTICAL_CALIBER_OPTIONS)}）、使用量、单位（t或m3）
- 氮气：统计口径、使用量、单位

请从提供的原始文本中提取上述字段，以JSON格式返回。只返回JSON。""",
    }


SECTION_PROMPTS = _build_section_prompts()


class LLMExtractor:
    """使用 LLM 从文本中提取结构化数据"""

    def __init__(self, section: int):
        self.section = section
        self.prompt = SECTION_PROMPTS.get(section, "")

    def extract(self, raw_text: str) -> Dict[str, Any]:
        """调用 LLM 提取数据"""
        if not raw_text.strip():
            return {}

        try:
            llm_client = get_llm_client()

            # 检查是否包含图片数据
            if "[IMAGE_DATA]" in raw_text:
                # 构建多模态消息
                messages = self._build_multimodal_message(raw_text)
                result = llm_client.chat(messages)
            else:
                # 纯文本消息
                messages = [
                    {"role": "system", "content": self.prompt},
                    {"role": "user", "content": f"原始文本：\n{raw_text[:8000]}"}
                ]
                result = llm_client.chat(messages)

            if isinstance(result, dict) and result.get("content"):
                answer = result["content"]
                logger.info(f"LLM原始输出 section={self.section}: {answer[:500]}")
                # 尝试解析 JSON
                parsed = self._parse_json(answer)
                logger.info(f"LLM解析结果 section={self.section}: {parsed}")
                return parsed
        except Exception as e:
            logger.error(f"LLM提取失败 section={self.section}: {e}")

        return {}

    def _build_multimodal_message(self, raw_text: str) -> List[Dict[str, Any]]:
        """构建多模态消息，处理图片和文本混合内容"""
        import re
        messages = [
            {"role": "system", "content": self.prompt},
            {"role": "user", "content": []}
        ]

        # 分割文本和图片数据
        parts = re.split(r'\[IMAGE_DATA\](.*?)\[/IMAGE_DATA\]', raw_text, flags=re.DOTALL)

        content = messages[1]["content"]

        for i, part in enumerate(parts):
            if i % 2 == 1:
                # 这是图片数据
                b64_data = part.split(" dimensions:")[0]
                content.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/jpeg;base64,{b64_data}"
                    }
                })
            elif part.strip() and part.strip() != "dimensions:":
                # 这是普通文本
                content.append({
                    "type": "text",
                    "text": part.strip()
                })

        # 过滤空内容
        messages[1]["content"] = [c for c in content if isinstance(c, dict) and (c.get("text") or c.get("image_url"))]

        return messages

    def _parse_json(self, text: str) -> Dict[str, Any]:
        """从 LLM 输出中解析 JSON"""
        import re
        import json

        # 尝试多种方式解析 JSON
        # 方式1：尝试直接解析整个文本（如果它是纯 JSON）
        text = text.strip()
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            pass

        # 方式2：尝试匹配 JSON 对象（支持嵌套大括号）
        # 找到第一个 { 和最后一个 } 之间的内容
        start = text.find('{')
        end = text.rfind('}')
        if start != -1 and end != -1 and end > start:
            json_str = text[start:end+1]
            try:
                return json.loads(json_str)
            except json.JSONDecodeError:
                pass

        # 方式3：尝试匹配数组格式的 JSON
        arr_start = text.find('[')
        arr_end = text.rfind(']')
        if arr_start != -1 and arr_end != -1 and arr_end > arr_start:
            json_str = text[arr_start:arr_end+1]
            try:
                return json.loads(json_str)
            except json.JSONDecodeError:
                pass

        logger.warning(f"无法解析 LLM 输出为 JSON: {text[:200]}")
        return {}

class FileExtractAgent:
    """文件提取 Agent - 组合 TextExtractor 和 LLMExtractor。

    Phase 0: 添加此缺失类以解除 import 阻塞。
    接口: FileExtractAgent(section).process(content, filename=...) -> dict
    """

    SUPPORTED_EXTENSIONS = {
        '.xlsx', '.xls', '.pdf', '.docx', '.doc', '.pptx', '.md',
        '.png', '.jpg', '.jpeg',
    }

    def __init__(self, section: int):
        self.section = section

    def process(self, content: bytes, filename: str = "") -> Dict[str, Any]:
        """从文件字节内容中提取对应 section 的数据。"""
        ext = os.path.splitext(filename)[1].lower() if filename else ''
        if ext not in self.SUPPORTED_EXTENSIONS:
            return {
                "status": "failed",
                "error": f"不支持的文件类型: {ext or '(无)'}",
                "data": {},
            }
        try:
            text = TextExtractor.extract_from_bytes(content, filename)
        except Exception as e:
            logger.error(f"文本提取失败: filename={filename}, error: {e}", exc_info=True)
            return {"status": "failed", "error": f"文本提取失败: {e}", "data": {}}
        if not text:
            return {"status": "failed", "error": "文件内容为空", "data": {}}
        try:
            data = LLMExtractor(self.section).extract(text)
        except Exception as e:
            logger.error(f"LLM 提取失败: section={self.section}, error: {e}", exc_info=True)
            return {"status": "failed", "error": f"AI 提取失败: {e}", "data": {}}
        return {"status": "completed", "data": data, "error": None}


class ExtractorsFactory:
    """提取器工厂 - 为每个部分创建对应的提取器"""

    @staticmethod
    def get_extractor(section: int) -> FileExtractAgent:
        """获取指定部分的提取器"""
        return FileExtractAgent(section)

    @staticmethod
    def get_all_extractors() -> Dict[int, FileExtractAgent]:
        """获取所有9个部分的提取器"""
        return {i: FileExtractAgent(i) for i in range(1, 10)}


